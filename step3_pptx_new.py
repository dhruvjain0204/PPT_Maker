"""
Step 3: PPTX Generation (New Implementation)
Generate PowerPoint presentation from structured question data using python-pptx.

This implementation follows python-pptx best practices:
- Proper use of text frames, paragraphs, and runs
- Clean separation of concerns
- Better error handling
- Improved table rendering
- Professional formatting

Reference: https://python-pptx.readthedocs.io/en/latest/
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class PPTXGenerator:
    """
    PowerPoint presentation generator for practice questions.
    
    Handles creation of slides with questions, answers, tables, and formatting.
    """
    
    # Slide dimensions (standard PowerPoint widescreen 16:9)
    # Standard widescreen: 13.333" × 7.5" (33.867 cm × 19.05 cm)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Text box dimensions and position based on Canva template
    # Desired Canva values: Width=1773.3px, Height=510.7px, X=73.6px, Y=162.3px
    # Current Canva values (from imported PPTX): Width=1745.8px, Height=482.3px, X=88.1px, Y=176.8px
    # Need to calculate correction factors to achieve desired Canva values
    
    # Canva canvas for widescreen is typically 1920px × 1080px
    CANVA_CANVAS_WIDTH_PX = 1920  # Standard Canva widescreen width
    CANVA_CANVAS_HEIGHT_PX = 1080  # Standard Canva widescreen height
    
    # Convert Canva pixel positions to PowerPoint inches
    # Scale factor: PowerPoint width (13.333") / Canva width (1920px)
    SCALE_X = 13.333 / CANVA_CANVAS_WIDTH_PX
    SCALE_Y = 7.5 / CANVA_CANVAS_HEIGHT_PX
    
    # Calculate correction ratios based on actual vs desired Canva values
    # When we set X based on 73.6px, Canva shows 88.1px
    # Correction ratio: desired / actual = 73.6 / 88.1 = 0.835
    CORRECTION_X = 73.6 / 88.1  # ~0.835
    CORRECTION_Y = 162.3 / 176.8  # ~0.918
    CORRECTION_WIDTH = 1773.3 / 1745.8  # ~1.016
    CORRECTION_HEIGHT = 510.7 / 482.3  # ~1.059
    
    # Apply corrections to get PowerPoint values that will result in desired Canva values
    # Text box position and size (corrected to achieve desired Canva values)
    TEXTBOX_LEFT = Inches(73.6 * SCALE_X * CORRECTION_X)      # Will result in Canva X=73.6px
    TEXTBOX_TOP = Inches(162.3 * SCALE_Y * CORRECTION_Y)       # Will result in Canva Y=162.3px
    TEXTBOX_WIDTH = Inches(1773.3 * SCALE_X * CORRECTION_WIDTH)  # Will result in Canva Width=1773.3px
    TEXTBOX_HEIGHT = Inches(510.7 * SCALE_Y * CORRECTION_HEIGHT)  # Will result in Canva Height=510.7px
    
    # Content area margins (for fallback/other elements)
    MARGIN_LEFT = Inches(0.75)
    MARGIN_TOP = Inches(0.5)
    CONTENT_WIDTH = Inches(12.0)
    CONTENT_HEIGHT = Inches(6.5)
    
    # Font settings
    FONT_NAME = "Frankfurter Medium"  # Font as per user requirement
    
    # Font size calculation for Canva:
    # Current: 32pt PPT → 48pt Canva (ratio = 48/32 = 1.5)
    # Current: 28pt PPT → 42pt Canva (ratio = 42/28 = 1.5)
    # Desired: 38pt in Canva for all slides
    # Required PPT size: 38 / 1.5 = 25.33pt (round to 25pt for cleaner value)
    FONT_SIZE_UNIFIED = Pt(25)  # Will result in ~38pt in Canva for all slides
    FONT_SIZE_QUESTION = FONT_SIZE_UNIFIED
    FONT_SIZE_ANSWER = FONT_SIZE_UNIFIED
    FONT_SIZE_PASSAGE = FONT_SIZE_UNIFIED
    FONT_SIZE_TABLE = Pt(18)  # Table font size (increased to 18pt)
    
    def __init__(self):
        """Initialize a new presentation."""
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        self.answer_counter = 1
    
    def create_blank_slide(self):
        """
        Create a blank slide with no placeholders.
        
        Returns:
            Slide object
        """
        # Use blank layout (index 6 is typically blank)
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        return slide
    
    def add_text_box(self, slide, text: str, left: Inches, top: Inches, 
                     width: Inches, height: Inches = None, 
                     font_size: Pt = None, bold: bool = False,
                     alignment: PP_ALIGN = PP_ALIGN.LEFT):
        """
        Add a text box to a slide with proper formatting.
        
        Args:
            slide: Slide object
            text: Text content
            left: Left position in Inches
            top: Top position in Inches
            width: Width in Inches
            height: Height in Inches (None for auto-sizing)
            font_size: Font size in Points
            bold: Whether text should be bold
            alignment: Text alignment
            
        Returns:
            TextFrame object
        """
        # Use maximum available height if not specified
        if height is None:
            height = self.CONTENT_HEIGHT
        
        # Create text box shape
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_bottom = Inches(0.15)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.auto_size = None  # Disable auto-size for manual control
        
        # Clear default paragraph
        text_frame.clear()
        
        # Handle multi-line text - split by newlines and create paragraphs
        lines = text.split('\n')
        
        for i, line in enumerate(lines):
            if i > 0:
                # Add new paragraph for each line break
                paragraph = text_frame.add_paragraph()
            else:
                paragraph = text_frame.paragraphs[0]
            
            paragraph.alignment = alignment
            paragraph.space_after = Pt(8) if i < len(lines) - 1 else Pt(0)
            paragraph.space_before = Pt(0)
            
            # Check if line starts with Q1, Q2, Ans1, Ans2, etc. for bold formatting
            line_stripped = line.strip()
            is_numbered = (line_stripped.startswith('Q') and len(line_stripped) > 1 and line_stripped[1:2].isdigit()) or \
                         (line_stripped.startswith('Ans') and len(line_stripped) > 3 and line_stripped[3:4].isdigit())
            
            if is_numbered:
                # Check if line already has a dot after the number (e.g., "Q1. text" or "Ans1. text")
                # Find where the number ends
                if line_stripped.startswith('Q'):
                    # Q1, Q2, Q10, etc.
                    num_end = 1
                    while num_end < len(line_stripped) and line_stripped[num_end].isdigit():
                        num_end += 1
                    number_part = line_stripped[:num_end]
                    rest = line_stripped[num_end:].strip()
                else:  # Ans
                    # Ans1, Ans2, Ans10, etc.
                    num_end = 3
                    while num_end < len(line_stripped) and line_stripped[num_end].isdigit():
                        num_end += 1
                    number_part = line_stripped[:num_end]
                    rest = line_stripped[num_end:].strip()
                
                # Check if rest starts with a dot (already has dot)
                if rest.startswith('.'):
                    # Already has dot, remove it and any spaces after it
                    text_part = rest.lstrip('.').lstrip().strip()
                    separator = ". " if text_part else "."
                elif rest.startswith(' '):
                    # Has space but no dot - add dot before space
                    text_part = rest.lstrip()
                    separator = ". "
                else:
                    # No dot or space, add dot and space
                    text_part = rest
                    separator = ". "
                
                # Add bold number run
                run1 = paragraph.add_run()
                run1.text = number_part + separator
                run1.font.name = self.FONT_NAME
                run1.font.size = font_size or self.FONT_SIZE_QUESTION
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(0, 0, 0)
                
                # Add regular text run
                if text_part:
                    run2 = paragraph.add_run()
                    run2.text = text_part
                    run2.font.name = self.FONT_NAME
                    run2.font.size = font_size or self.FONT_SIZE_QUESTION
                    run2.font.bold = bold
                    run2.font.color.rgb = RGBColor(0, 0, 0)
            else:
                # Regular line - add single run
                run = paragraph.add_run()
                run.text = line if line.strip() else " "  # Preserve empty lines
                run.font.name = self.FONT_NAME
                run.font.size = font_size or self.FONT_SIZE_QUESTION
                run.font.bold = bold
                run.font.color.rgb = RGBColor(0, 0, 0)
        
        return text_frame
    
    def add_formatted_text(self, text_frame, text: str, font_size: Pt = None, 
                          bold: bool = False, color: RGBColor = None):
        """
        Add formatted text to an existing text frame.
        
        Args:
            text_frame: TextFrame object
            text: Text to add
            font_size: Font size in Points
            bold: Whether text should be bold
            color: Text color (default: black)
        """
        if color is None:
            color = RGBColor(0, 0, 0)
        
        # Split by newlines to preserve line breaks
        lines = text.split('\n')
        
        for i, line in enumerate(lines):
            # Add new paragraph for each line
            paragraph = text_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8) if i < len(lines) - 1 else Pt(0)
            paragraph.space_before = Pt(0)
            
            # Add run
            run = paragraph.add_run()
            run.text = line if line.strip() else " "
            run.font.name = self.FONT_NAME
            run.font.size = font_size or self.FONT_SIZE_QUESTION
            run.font.bold = bold
            run.font.color.rgb = color
    
    def add_table(self, slide, table_data: dict, left: Inches, top: Inches,
                  width: Inches = None, height: Inches = None):
        """
        Add a table to a slide.
        
        Args:
            slide: Slide object
            table_data: Dictionary with 'headers' and 'rows'
            left: Left position in Inches
            top: Top position in Inches
            width: Table width in Inches (default: CONTENT_WIDTH)
            height: Table height in Inches (auto-calculated if None)
            
        Returns:
            Table object
        """
        if not table_data or not table_data.get('headers'):
            return None
        
        headers = table_data['headers']
        rows = table_data['rows']
        
        if width is None:
            width = self.CONTENT_WIDTH
        
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header row
        
        # Calculate height if not provided
        if height is None:
            # Estimate: ~0.4 inches per row, max 4 inches
            estimated_height = min(num_rows * 0.4, 4.0)
            height = Inches(estimated_height)
        
        # Create table shape
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths proportionally
        col_width = int((width.inches / num_cols) * 914400)  # Convert to EMU
        for col_idx in range(num_cols):
            table.columns[col_idx].width = col_width
        
        # Format header row
        header_row = table.rows[0]
        for col_idx, header_text in enumerate(headers):
            cell = header_row.cells[col_idx]
            cell.text_frame.clear()
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            run = paragraph.add_run()
            run.text = str(header_text) if header_text else ""
            run.font.name = self.FONT_NAME
            run.font.size = self.FONT_SIZE_TABLE
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Set cell fill color (light gray for headers)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        # Format data rows
        for row_idx, row_data in enumerate(rows, 1):
            table_row = table.rows[row_idx]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table_row.cells[col_idx]
                    cell.text_frame.clear()
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    run = paragraph.add_run()
                    run.text = str(cell_text) if cell_text is not None else ""
                    run.font.name = self.FONT_NAME
                    run.font.size = self.FONT_SIZE_TABLE
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        return table
    
    def create_question_slide(self, question_data: dict, question_number: str):
        """
        Create a question slide with question text, options, table, and diagram.
        
        Args:
            question_data: Content dictionary from parsed questions
            question_number: Question number (e.g., "Q1")
        """
        slide = self.create_blank_slide()
        
        # Use Canva template dimensions for text box
        left = self.TEXTBOX_LEFT
        top = self.TEXTBOX_TOP
        width = self.TEXTBOX_WIDTH
        
        # Build full question text with all components
        question_text = question_data.get('question_text', '')
        full_text_parts = [f"{question_number}. {question_text}"]
        
        # Add options if present (without bullets for multiple choice)
        options = question_data.get('options', [])
        if options:
            options_text = "\n".join(opt for opt in options)
            full_text_parts.append(options_text)
        
        # Add diagram note if present (just show [Diagram] instead of full description)
        diagram = question_data.get('diagram_description')
        if diagram:
            full_text_parts.append("\n[Diagram]")
        
        # Combine all text
        full_question = "\n\n".join(full_text_parts)
        
        # Check if we have a table
        table = question_data.get('table')
        has_table = table and table.get('headers')
        
        # Determine if table is actually just options in table format
        # If options exist and table exists, check if table represents the same options
        is_options_table = False
        if has_table and options:
            table_rows = table.get('rows', [])
            # Check if table rows match options pattern (e.g., first column has "(a)", "(b)", etc.)
            if table_rows and len(table_rows) == len(options):
                # Check if first column of each row matches option labels
                matches = 0
                for i, row in enumerate(table_rows):
                    if row and len(row) > 0:
                        first_cell = str(row[0]).strip()
                        # Check if it matches option label pattern like "(a)", "(b)", etc.
                        if first_cell.startswith('(') and len(first_cell) >= 3 and first_cell[1:2].isalpha() and first_cell.endswith(')'):
                            # Check if it matches the corresponding option
                            option_label = options[i].strip()[:3] if len(options[i]) >= 3 else ""
                            if first_cell.lower() == option_label.lower():
                                matches += 1
                
                # If most rows match, it's likely an options table
                if matches >= len(options) * 0.75:  # 75% match threshold
                    is_options_table = True
        
        # Only show table if it's NOT an options table (normal table)
        show_table = has_table and not is_options_table
        
        # Calculate available height for text
        if show_table:
            # Reserve space for table below text box
            text_height = self.TEXTBOX_HEIGHT - Inches(1.0)
        else:
            text_height = self.TEXTBOX_HEIGHT
        
        # Create text frame for question (with all content) using Canva dimensions
        text_frame = self.add_text_box(
            slide, full_question, left, top, width,
            height=text_height,
            font_size=self.FONT_SIZE_QUESTION, bold=False
        )
        
        # Add table below text if present and it's a normal table (not options)
        if show_table:
            table_top = top + text_height + Inches(0.2)
            table_width = min(width, self.CONTENT_WIDTH)  # Use appropriate width
            self.add_table(slide, table, left, table_top, width=table_width)
    
    def create_answer_slide(self, answer_text: str):
        """
        Create an answer slide with answer text.
        
        Args:
            answer_text: Answer content
        """
        slide = self.create_blank_slide()
        
        answer_number = f"Ans{self.answer_counter}"
        self.answer_counter += 1
        
        full_answer = f"{answer_number}. {answer_text}"
        
        # Use Canva template dimensions for answer slide too
        self.add_text_box(
            slide, full_answer,
            self.TEXTBOX_LEFT, self.TEXTBOX_TOP,
            self.TEXTBOX_WIDTH,
            height=self.TEXTBOX_HEIGHT,
            font_size=self.FONT_SIZE_ANSWER,
            bold=False
        )
    
    def create_passage_slide(self, passage_text: str):
        """
        Create a passage slide with passage text.
        
        Args:
            passage_text: Passage content
        """
        slide = self.create_blank_slide()
        
        # Use Canva template dimensions for passage slide
        self.add_text_box(
            slide, passage_text,
            self.TEXTBOX_LEFT, self.TEXTBOX_TOP,
            self.TEXTBOX_WIDTH,
            height=self.TEXTBOX_HEIGHT,
            font_size=self.FONT_SIZE_PASSAGE,
            bold=False
        )
    
    def generate(self, questions: list, output_file: str, include_answers: bool = True):
        """
        Generate PowerPoint presentation from parsed questions.
        
        Args:
            questions: List of question dictionaries
            output_file: Path to output PPTX file
            include_answers: Whether to include answer slides (default: True)
        """
        print("[INFO] Creating PowerPoint presentation...")
        if not include_answers:
            print("[INFO] Answer slides will be excluded from the presentation")
        
        # Process each question
        for q in questions:
            question_number = q.get('question_number', 'Q?')
            slides = q.get('slides', [])
            
            # Process each slide for this question
            for slide_data in slides:
                slide_type = slide_data.get('slide_type')
                content = slide_data.get('content', {})
                
                if slide_type == 'passage':
                    passage = content.get('passage', '')
                    self.create_passage_slide(passage)
                    
                elif slide_type == 'question':
                    self.create_question_slide(content, question_number)
                    
                elif slide_type == 'answer':
                    # Skip answer slides if include_answers is False
                    if include_answers:
                        answer_text = content.get('answer_text', '')
                        self.create_answer_slide(answer_text)
                    # If include_answers is False, skip this slide entirely
        
        # Save presentation
        Path(output_file).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_file)
        print(f"[OK] Presentation saved to: {output_file}")


def load_parsed_questions(json_file: str) -> list:
    """
    Load parsed questions from JSON file.
    
    Args:
        json_file: Path to parsed questions JSON file
        
    Returns:
        List of question dictionaries, or None if error
    """
    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            questions = json.load(f)
        print(f"[OK] Loaded {len(questions)} questions from {json_file}")
        return questions
    except FileNotFoundError:
        print(f"[ERROR] File not found: {json_file}")
        print("[INFO] Please run Step 2 first to parse questions")
        return None
    except json.JSONDecodeError as e:
        print(f"[ERROR] Failed to parse JSON: {e}")
        return None


def main():
    """Main execution function."""
    import sys
    
    # Check for --no-answers flag
    include_answers = True
    if len(sys.argv) > 1 and '--no-answers' in sys.argv:
        include_answers = False
        print("[INFO] Answer slides will be excluded from the presentation")
    
    # Read PDF name from Step 1
    pdf_name = "Adobe-Scan-12-Dec-2025"  # default fallback
    pdf_name_file = Path("output/current_pdf_name.txt")
    if pdf_name_file.exists():
        with open(pdf_name_file, 'r', encoding='utf-8') as f:
            pdf_name = f.read().strip()
        print(f"[INFO] Using PDF name from Step 1: {pdf_name}")
    else:
        print(f"[INFO] Using default PDF name: {pdf_name}")
    
    # Input file from Step 2 (named after PDF)
    input_file = f"output/parsed_questions_{pdf_name}.json"
    
    # Print header
    print("=" * 60)
    print("STEP 3: PPTX GENERATION (New Implementation)")
    print("=" * 60)
    print()
    
    # Load parsed questions
    questions = load_parsed_questions(input_file)
    
    if not questions:
        print("[ERROR] Failed to load parsed questions")
        return
    
    # Count slides (excluding answers if include_answers is False)
    if include_answers:
        total_slides = sum(len(q.get('slides', [])) for q in questions)
    else:
        # Count only non-answer slides
        total_slides = 0
        for q in questions:
            for slide in q.get('slides', []):
                if slide.get('slide_type') != 'answer':
                    total_slides += 1
    
    answer_status = "with answers" if include_answers else "without answers"
    print(f"[INFO] Will generate {total_slides} slides from {len(questions)} questions ({answer_status})")
    
    # Create PPTs folder if it doesn't exist
    ppt_folder = Path("PPTs")
    ppt_folder.mkdir(exist_ok=True)
    
    # Generate PPTX (named after PDF, saved to PPTs folder)
    generator = PPTXGenerator()
    output_file = f"PPTs/{pdf_name}.pptx"
    generator.generate(questions, output_file, include_answers=include_answers)
    
    # Verify output
    try:
        verify_prs = Presentation(output_file)
        num_slides = len(verify_prs.slides)
        print(f"\n[SUMMARY]")
        print(f"Total questions: {len(questions)}")
        print(f"Total slides generated: {num_slides}")
        if not include_answers:
            print(f"Note: Answer slides were excluded")
        print(f"Output file: {output_file}")
        print(f"\n[SUCCESS] Presentation generated successfully!")
        print(f"[INFO] Open the PPTX file to review the results")
    except Exception as e:
        print(f"[WARNING] Could not verify output file: {e}")


if __name__ == '__main__':
    main()

